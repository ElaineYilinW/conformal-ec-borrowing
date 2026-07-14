#!/usr/bin/env python3
# Build the presentation deck (Google-Slides-importable .pptx) from the figure
# and formula PNGs. Run from the CC_ppt/ folder:  python3 build_ppt.py
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
def F(name): return os.path.join(HERE, "figs", name + ".png")
def E(name): return os.path.join(HERE, "eq",   name + ".png")

SW, SH = 13.333, 7.5
INK   = RGBColor(0x1F,0x2A,0x44); BLUE = RGBColor(0x2F,0x7D,0xC8)
GREEN = RGBColor(0x3E,0x7A,0x1E); AMBER= RGBColor(0xB5,0x6A,0x10)
RED   = RGBColor(0xB0,0x2E,0x2E); BODY = RGBColor(0x33,0x37,0x40)
MUT   = RGBColor(0x6B,0x72,0x80); WHITE= RGBColor(0xFF,0xFF,0xFF)
LIGHT = RGBColor(0xEE,0xF2,0xF8); PANEL= RGBColor(0xF5,0xF7,0xFA)

prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]
FONT = "Arial"

def newslide():
    return prs.slides.add_slide(BLANK)

def rect(s,l,t,w,h,fill,line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l),Inches(t),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def tbox(s,l,t,w,h,anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf

def run(p, txt, size, color=BODY, bold=False, italic=False):
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.italic = italic; r.font.name = FONT
    return r

def title(s, ttl, kicker=None):
    rect(s, 0, 0, SW, 0.14, BLUE)
    tf = tbox(s, 0.55, 0.30, SW-1.1, 1.02)
    if kicker:
        p = tf.paragraphs[0]; run(p, kicker.upper(), 12.5, BLUE, bold=True); p.space_after = Pt(2)
        p2 = tf.add_paragraph()
    else:
        p2 = tf.paragraphs[0]
    run(p2, ttl, 27, INK, bold=True)
    rect(s, 0.57, 1.34, 2.3, 0.045, BLUE)

def bullets(s, items, l, t, w, h, size=15, gap=8, anchor=MSO_ANCHOR.TOP):
    tf = tbox(s, l, t, w, h, anchor); first=True
    for it in items:
        lvl, txt, kind, col, bold = 0, it, "b", BODY, False
        if isinstance(it, tuple):
            kind = it[0]; txt = it[1]
            col  = it[2] if len(it) > 2 else BODY
            bold = it[3] if len(it) > 3 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        if kind == "h":            # section header line (no bullet)
            run(p, txt, size+1.5, BLUE, bold=True); p.space_after = Pt(gap+1); p.space_before = Pt(4)
        elif kind == "n":          # plain line, no bullet
            run(p, txt, size, col, bold=bold); p.space_after = Pt(gap)
        elif kind == "s":          # sub-bullet
            p.level = 1; run(p, "–  ", size-1, MUT); run(p, txt, size-1, col, bold=bold); p.space_after = Pt(gap-2)
        else:                       # normal bullet
            run(p, "▪  ", size, BLUE, bold=True); run(p, txt, size, col, bold=bold); p.space_after = Pt(gap)
    return tf

def img(s, path, l, t, w, h, halign="center", valign="middle", bg=False):
    iw, ih = Image.open(path).size; ar = iw/ih; bar = w/h
    if ar >= bar: nw, nh = w, w/ar
    else:         nw, nh = h*ar, h
    lx = l + (w-nw)/2 if halign=="center" else (l if halign=="left" else l+(w-nw))
    ty = t + (h-nh)/2 if valign=="middle" else (t if valign=="top" else t+(h-nh))
    if bg: rect(s, lx-0.12, ty-0.10, nw+0.24, nh+0.20, PANEL)
    s.shapes.add_picture(path, Inches(lx), Inches(ty), Inches(nw), Inches(nh))
    return (lx, ty, nw, nh)

def pagenum(s, n):
    tf = tbox(s, SW-1.2, SH-0.44, 1.0, 0.35)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; run(p, str(n), 11, MUT)

def section(kicker, ttl, subtitle=None):
    s = newslide(); rect(s, 0,0, SW, SH, INK); rect(s, 0, 3.02, SW, 0.05, BLUE)
    tf = tbox(s, 1.1, 2.35, SW-2.2, 2.2)
    p = tf.paragraphs[0]; run(p, kicker.upper(), 17, BLUE, bold=True); p.space_after = Pt(6)
    p2 = tf.add_paragraph(); run(p2, ttl, 40, WHITE, bold=True)
    if subtitle:
        p3 = tf.add_paragraph(); run(p3, subtitle, 16, RGBColor(0xC9,0xD4,0xE4), italic=True); p3.space_before = Pt(8)
    return s

N = [0]
def page(s):
    N[0]+=1; pagenum(s, N[0])

# ============================================================ TITLE
s = newslide(); rect(s,0,0,SW,SH,INK); rect(s,0,4.5,SW,0.06,BLUE)
tf = tbox(s, 1.0, 1.7, SW-2.0, 2.6)
p=tf.paragraphs[0]; run(p,"Safe Borrowing of External Controls", 34, WHITE, bold=True); p.space_after=Pt(6)
p2=tf.add_paragraph(); run(p2,"via Covariate-Shift-Adjusted Longitudinal Residual-Rank Screening", 22, RGBColor(0x8F,0xB0,0xDE));
tf2 = tbox(s, 1.0, 4.75, SW-2.0, 2.0)
for line,c,sz in [("A conformal screen for hybrid controlled trials",RGBColor(0xC9,0xD4,0xE4),16),
                  ("Keep the compatible external controls · drop the incompatible ones · borrow safely",RGBColor(0x9A,0xA7,0xBC),14)]:
    pp=tf2.add_paragraph(); run(pp,line,sz,c,italic=(sz==14)); pp.space_after=Pt(4)

# ============================================================ PART I
section("Part I", "The Problem & Background", "Why borrow external controls, and when is it dangerous?")

# --- Motivation
s=newslide(); title(s,"Hybrid trials: borrowing external controls","Scientific motivation")
bullets(s,[
 ("h","The setting"),
 "A randomized trial (RCT) gives unbiased treatment-effect estimates — but its control arm is often small, so precision and power are limited.",
 "Registries, historical trial arms, and electronic health records offer many external controls (ECs) that could sharpen the control comparison.",
 ("h","The danger"),
 "Borrowing ALL of them is risky: if a hidden subset is incompatible (their control-outcome law really differs), the borrowed control mean is biased and confidence-interval coverage collapses.",
 "RCT-only is always valid but discards the external information — precision left on the table.",
 ("n","Goal: screen the external controls for outcome compatibility — keep the compatible ones, drop the rest — then borrow. “Screen, then borrow.”",INK,True),
], 0.55,1.55,6.55,5.5, size=14.5)
img(s, F("fig_design"), 7.25, 1.7, 5.6, 4.9)
page(s)

# --- Why longitudinal
s=newslide(); title(s,"Why longitudinal outcomes matter","Scientific motivation")
bullets(s,[
 "Outcomes are trajectories over visits, not a single number. The estimand is a prespecified linear contrast of the trajectory (formula at right):",
 ("s","examples of the contrast c: final visit · final-minus-baseline · visit average · slope · area-under-the-curve."),
 "Longitudinal data reveal the temporal PATTERN of response: two sources can agree at the last visit yet differ in early decline, slope, curvature, or recovery.",
 "Such discrepancies can matter clinically — and can expose incompatibility that a single scalar endpoint cannot see.",
 "It is also harder to model: within-subject correlation, irregular visit timing, missingness, dropout.",
 ("n","Key takeaway: a single endpoint can hide incompatibility that the full trajectory reveals — so the screen compares whole trajectories, not one number.",INK,True),
], 0.55,1.5,6.5,5.4, size=14)
img(s, F("fig_long"), 7.1, 1.65, 5.75, 3.2)
img(s, E("eq_estimand"), 7.35, 5.45, 5.2, 0.8, bg=True)
page(s)

# --- Why screen-then-borrow
s=newslide(); title(s,"Real-world external controls differ in many ways","Scientific motivation")
bullets(s,[
 ("h","Ways compatible-looking ECs can actually differ"),
 "Baseline disease severity and prognostic factors;",
 "Visit schedules and observation frequency;",
 "Endpoint collection, measurement error, and adjudication;",
 "Missingness and dropout mechanisms;",
 "Calendar time, standard of care, concomitant treatment;",
 "Conditional outcome means, variances, correlations, or tails.",
], 0.55,1.55,6.2,5.4, size=15)
bullets(s,[
 ("h","The consequence"),
 ("n","Full borrowing improves precision ONLY if the ECs are compatible; otherwise it injects bias.",BODY,False),
 ("n","So the scientific question is not WHETHER to use external data, but WHICH external subjects carry credible information for the trial population.",INK,True),
 ("n","Our answer: screen at the OUTCOME level (residual trajectories) — complementary to design-stage covariate balancing.",BODY,False),
], 7.05,1.55,5.8,5.4, size=15)
page(s)

# --- Identification / notation
s=newslide(); title(s,"Identification: notation and the target","Identification assumptions & diagnostic targets")
bullets(s,[
 ("h","Notation"),
 "S = 1 the randomized trial, S = 0 the external source; within the trial A = 1 treated, A = 0 control.",
 "Screening compares randomized controls (S=1, A=0) with external controls (S=0), both observed under control.",
 "X = baseline covariates measured comparably in both sources.",
 ("h","The causal target"),
 "A prespecified contrast of the trajectory, in the randomized-trial population:",
], 0.55,1.55,7.0,3.6, size=15)
img(s, E("eq_estimand"), 0.7, 4.65, 6.6, 0.8, halign="left", bg=True)
bullets(s,[
 ("n","Central discipline of this work:",INK,True),
 ("s","separate the causal IDENTIFICATION assumptions (below) from the weaker object a chosen score actually SCREENS."),
 ("s","The two are not the same — conflating them is a common error."),
], 0.7,5.65,12.0,1.4, size=14)
page(s)

# --- Assumptions (support + compatibility)
s=newslide(); title(s,"Assumptions","Identification assumptions & diagnostic targets")
bullets(s,[
 ("h","Assumption 0 — covariate overlap (support)"),
 ("n","The trial's covariate support must be covered by the external source. The covariate DISTRIBUTIONS may still differ — that is covariate shift, which is allowed and corrected by weighting (Part II).",BODY,False),
], 0.55,1.5,7.05,1.55, size=14)
img(s, E("eq_overlap"), 0.65, 2.95, 6.9, 0.6, halign="left", bg=True)
bullets(s,[
 ("h","Compatibility of the control outcome law (nested, weakest → strongest)"),
 ("s","ME-c: the mean of the ONE estimand contrast is equal — all that IDENTIFICATION needs."),
 ("s","VME: the full mean vector is equal at every visit."),
 ("s","DE: the full conditional distribution is equal (what ordinary conformal validity nominally asks)."),
], 0.55,3.8,7.1,2.4, size=13.5)
img(s, E("eq_nest"), 0.7, 6.3, 6.6, 0.55, halign="left")
img(s, F("fig_compat"), 7.75, 1.65, 5.15, 5.2)
page(s)

# --- Existing design-stage screen (propensity) — BEFORE the residual-rank screen
s=newslide(); title(s,"Existing tool: design-stage propensity-score trimming","Existing design-stage screen")
bullets(s,[
 ("h","What it does"),
 "The source propensity score e(X) = Pr(S=1 | X); matching / weighting / trimming balance measured baseline covariates.",
 "Extreme e(X) flags regions where one source has little density relative to the other:",
], 0.55,1.55,7.0,2.1, size=15)
img(s, E("eq_ps"), 0.7, 3.35, 7.2, 0.75, halign="left", bg=True)
bullets(s,[
 ("h","What it does NOT do"),
 "Propensity trimming targets covariate OVERLAP / positivity only (Assumption 0).",
 "It does not assess whether the control OUTCOME law Y(0) | X is actually compatible across sources.",
 ("n","Complementary, not competing:",INK,True),
 ("s","propensity trimming  →  X-overlap;   residual-rank screening  →  outcome-compatibility through the score V."),
 ("s","Units outside the X-support should not be rescued by an outcome screen; units WITH good X-support may still have incompatible trajectories — that is what our screen catches next."),
], 0.55,4.35,12.3,2.7, size=14.5)
page(s)

# --- What a residual-rank screen actually checks
s=newslide(); title(s,"What a residual-rank screen actually checks","Identification assumptions & diagnostic targets")
bullets(s,[
 ("h","The object actually screened"),
 "A residual-rank screen tests the distribution of the nonconformity SCORE V — an object as strong as ME-c or VME depending on the score, but NEVER the full law DE.",
 "So it controls borrowing bias for the estimand without asserting the (essentially unfalsifiable) full-law equality: it sits between what identification NEEDS (ME-c) and the too-strong DE.",
 ("h","The score decides what you can see"),
 "WHICH score you rank on determines which departures the screen detects: a score aligned to the drift is powerful, a score blind to it has no power (we define and choose the score in Part II).",
 ("n","Design target: choose the score so that “looks compatible on the score” implies “compatible for your estimand (ME-c).”",INK,True),
], 0.55,1.55,12.3,5.4, size=14.5)
page(s)

# ============================================================ PART II
section("Part II", "The Method, Step by Step", "From a broad-audience idea to the full formulation")

# --- Building blocks (fundamentals): nonconformity score + split conformal
s=newslide(); title(s,"Building blocks: nonconformity scores & split conformal","Basic idea for a broad audience")
bullets(s,[
 ("h","A nonconformity score"),
 ("n","A scalar V = s(Y, X; model): how UNUSUAL a point is given the fitted model (e.g. a standardized residual; large |V| = far from expected).",BODY,False),
 ("h","Split conformal"),
 ("n","Fit on a TRAINING fold; score a held-out CALIBRATION fold; rank a new point's score among those calibration scores.",BODY,False),
 ("h","Why it is valid — and robust to a wrong model"),
 ("n","Under compatibility the new point is EXCHANGEABLE with the calibration controls → uniform rank → valid p-value, for ANY score.",BODY,False),
 ("n","Even if the model is MISSPECIFIED, controls and compatible ECs share the SAME score pattern → ranks unchanged → still valid. Misspecification costs POWER, not validity.",GREEN,False),
], 0.55,1.45,12.35,3.2, size=12)
img(s, F("fig_conformal_intro"), 0.55, 4.72, 12.2, 2.4)
page(s)

# --- Core idea applied to ECs (TWO-SIDED keep/drop)
s=newslide(); title(s,"The core idea, applied to the external controls","Basic idea for a broad audience")
bullets(s,[
 "For each external control, compute the same score and RANK it among the randomized-control scores — a conformal p-value:",
], 0.55,1.55,6.6,1.15, size=15)
img(s, E("eq_pbasic"), 0.7, 2.6, 4.4, 1.15, halign="left", bg=True)
bullets(s,[
 ("h","Keep the middle, drop BOTH tails"),
 ("n","A score that is extreme in EITHER direction — unusually high OR unusually low — signals incompatibility. So drop an EC whose p-value is too small OR too large, and keep the central band:",RED,False),
], 0.55,4.0,6.7,1.85, size=13.5)
img(s, E("eq_keepband"), 0.62, 5.7, 6.95, 0.6, halign="left", bg=True)
bullets(s,[
 ("n","The fitted model supplies the ORDERING; the randomized controls supply how extreme that ordering is.",INK,True),
], 0.55,6.45,6.7,0.7, size=13)
img(s, F("fig_rank"), 7.05, 1.6, 5.85, 5.2)
page(s)

# --- Why rank > residual (1)
s=newslide(); title(s,"Why rank, not a raw residual — (1) distribution-free","Why rank calibration is preferable")
bullets(s,[
 "After cross-fitting, regard the working model as fixed. Under compatibility the control and EC scores are EXCHANGEABLE, so their conditional CDFs match:",
], 0.55,1.5,6.6,1.4, size=15)
img(s, E("eq_exch"), 0.6, 2.75, 6.6, 0.95, halign="left", bg=True)
bullets(s,[
 ("n","This holds for ANY fixed score — even one from a misspecified model.",INK,True),
 ("n","⇒ the rank is uniform, the p-value super-uniform (valid): no normality, no chi-square reference, no numeric threshold to trust.",BODY,False),
 ("h","Contrast with a fixed residual cutoff"),
 "A cutoff |r| > 2σ̂ must trust the model's SCALE σ̂. If the model understates the spread, the cutoff excludes many COMPATIBLE controls (right).",
], 0.55,3.85,6.6,3.1, size=14)
img(s, F("fig_rankres"), 7.05, 2.0, 5.85, 3.6)
page(s)

# --- Why rank > residual (2) misspec  (CORE)
s=newslide(); title(s,"Why rank — (2) robust to model misspecification","The core reason to rank, not threshold")
bullets(s,[
 ("h","The deeper reason to rank rather than threshold the residual value"),
 "A wrong working mean model shifts EVERY compatible residual — randomized controls AND compatible ECs alike — by the SAME bias b (they share the true conditional law, so the same misspecified score map applies to both).",
 ("n","A fixed cutoff reads that shared bias as incompatibility  →  excludes most compatible controls.",RED,False),
 ("n","Ranking against the controls: their quantile carries the same bias, so it MOVES with them  →  the shared error cancels  →  only genuinely different ECs are flagged.",GREEN,False),
 ("n","This is the conformal exchangeability guarantee — valid p-values under ANY working model, correct or not. Scale-robustness (previous slide) is a special case.",INK,True),
], 0.55,1.5,12.3,2.9, size=14.5)
img(s, F("fig_misspec"), 0.9, 4.5, 11.5, 2.5)
page(s)

# --- Full method overview
s=newslide(); title(s,"The full method: six steps","Proposed full method")
img(s, F("fig_pipeline"), 0.5, 1.5, 12.3, 3.0)
bullets(s,[
 ("n","LEARN what a normal control looks like:  (1) fit a control-mean model on RCT controls,  (2) residualize,  (3) reduce to an estimand-aligned score.",BODY,False),
 ("n","FLAG the abnormal ECs:  (4) rank each EC against the RCT controls with covariate-shift weighting,  (5) screen with a symmetric, adaptive threshold.",BODY,False),
 ("n","BORROW the rest:  (6) refit a transported AIPW on RCT controls + retained ECs  →  τ̂ with a confidence interval.",INK,True),
], 0.7,4.75,12.0,2.2, size=14.5)
page(s)

# --- Step 1-2
s=newslide(); title(s,"Steps 1–2: estimand and working models","Proposed full method")
bullets(s,[
 ("h","Step 1 — prespecify the estimand & harmonize the data"),
 "Choose the primary contrast c and target τ_c; two contrasts are used later so a given drift projects onto them differently:",
], 0.55,1.55,12.2,1.7, size=15)
img(s, E("eq_estimands"), 0.7, 3.2, 8.6, 0.7, halign="left", bg=True)
bullets(s,[
 ("s","Align eligibility, endpoint definitions, visit windows, baseline covariates, data-quality rules; check X-overlap."),
 ("h","Step 2 — fit working longitudinal models on the randomized controls"),
 "Fit the control mean (and scale / covariance) models by K-fold CROSS-FITTING, so in-sample residuals are not over-optimistic.",
 ("n","The model need only supply a useful ordering — rank calibration handles the rest, even under misspecification.",INK,True),
], 0.55,4.1,12.3,2.9, size=14.5)
page(s)

# --- Step 3 scores
s=newslide(); title(s,"Step 3: the longitudinal nonconformity score","Proposed full method")
bullets(s,[
 "Residualize the trajectory,  r = Y − m̂(X),  then reduce it to a single scalar score. Three options:",
], 0.55,1.5,12.2,1.0, size=15)
img(s, E("eq_scores"), 0.7, 2.4, 11.6, 0.85, halign="left", bg=True)
bullets(s,[
 ("n","Contrast score (aligned to the estimand c): most powerful when the drift is along c — global shifts, level changes.",BODY,False),
 ("n","Max standardized residual: sensitive to a spike at any single visit.",BODY,False),
 ("n","Common-reference Mahalanobis: aggregates the whole trajectory shape — best for ramps / curvature / transient dips.",BODY,False),
 ("h","Design principle"),
 ("n","Choose the score to match the ANTICIPATED drift shape. A contrast blind to the drift will not detect it (we show this in simulation).",INK,True),
], 0.55,3.55,12.3,3.4, size=15)
page(s)

# --- Step 4-5 weighting
s=newslide(); title(s,"Steps 4–5: covariate-shift weighting","Proposed full method")
bullets(s,[
 "Randomized and external covariate mixes differ, so we reweight the calibration ranks by the density ratio:",
], 0.55,1.5,6.6,1.0, size=14.5)
img(s, E("eq_omega"), 0.65, 2.35, 6.5, 0.62, halign="left", bg=True)
img(s, E("eq_wpval"), 0.65, 3.15, 5.6, 1.1, halign="left", bg=True)
bullets(s,[
 ("n","Transport identity — weighted RCT scores represent the compatible-EC score law:",BODY,False),
], 0.55,4.35,6.7,0.8, size=13.5)
img(s, E("eq_transport"), 0.65, 5.05, 6.6, 0.7, halign="left", bg=True)
bullets(s,[
 ("n","Without weighting, high-severity compatible ECs look “extreme” only because they are over-represented externally → false exclusions. Weighting removes exactly this (right).",INK,True),
], 0.55,5.95,6.7,1.1, size=13.5)
img(s, F("fig_shift"), 7.2, 1.75, 5.6, 5.0)
page(s)

# --- Step 6-7 CV+
s=newslide(); title(s,"Steps 6–7: from a one-time split to cross-fit (CV+)","Proposed full method")
bullets(s,[
 "A single train/calibrate split wastes half of the already-scarce randomized controls.",
 "CV+ scores each fold from a model fit on the other folds:",
], 0.55,1.55,6.5,1.6, size=15)
img(s, E("eq_cv"), 0.65, 3.15, 6.4, 0.72, halign="left", bg=True)
bullets(s,[
 ("n","Every randomized control eventually contributes a calibration rank  →  the WHOLE control sample is used.",BODY,False),
 ("n","This is the main efficiency lever of the method — we quantify the gain in the simulation (≈ +0.06 power at controlled type-I).",INK,True),
], 0.55,4.1,6.6,2.6, size=15)
img(s, F("fig_cv"), 7.05, 1.7, 5.85, 5.1)
page(s)

# --- Step 8
s=newslide(); title(s,"Step 8: conservative screening — symmetric + adaptive γ","Proposed full method")
bullets(s,[
 ("h","Symmetric trimming (keeps the borrowed mean unbiased)"),
 "Screen on |V| with a SINGLE threshold, trimming both tails equally:",
], 0.55,1.5,6.6,1.6, size=14.5)
img(s, E("eq_symscreen"), 0.65, 2.95, 6.5, 0.72, halign="left", bg=True)
bullets(s,[
 ("n","Trimming both tails equally leaves the retained control mean unmoved → the borrowed mean stays UNBIASED (right).",BODY,False),
 ("h","Adaptive threshold (no labels needed)"),
 "Choose γ per dataset to minimize an estimated MSE of the borrowing estimate (a de-biased CSB objective):",
], 0.55,3.9,6.6,2.2, size=14)
img(s, E("eq_adaptive"), 0.65, 6.0, 6.6, 0.85, halign="left", bg=True)
img(s, F("fig_trimming"), 7.15, 2.2, 5.7, 3.1)
bullets(s,[
 ("n","Small γ borrows more but risks leaking contamination; large γ is safer with less gain. Adaptive γ balances the two automatically.",INK,True),
], 7.15,5.5,5.7,1.4, size=13.5)
page(s)

# --- Step 9
s=newslide(); title(s,"Step 9: refit the borrowing estimator (transported AIPW)","Proposed full method")
bullets(s,[
 "On the RETAINED external controls, refit the source density, the outcome model, and the inverse-variance borrowing weight; form the transported AIPW control mean and combine it with the randomized-control mean:",
], 0.55,1.55,12.2,1.6, size=15)
img(s, E("eq_estimator"), 0.9, 3.3, 10.4, 1.0, halign="left", bg=True)
bullets(s,[
 ("n","λ = efficient inverse-variance weight: borrow LESS from noisier external controls → lower RMSE.",BODY,False),
 ("n","Doubly robust: consistent if either the outcome model or the source/weight model is correct.",BODY,False),
 ("n","Output: τ̂ with a full-pipeline bootstrap confidence interval that re-runs the entire screen on each resample (so selection uncertainty is propagated).",INK,True),
], 0.55,4.65,12.3,2.3, size=15)
page(s)

# ============================================================ PART III
section("Part III", "Simulation", "A transparent longitudinal data-generating process, and the results")

# --- Sim setup
s=newslide(); title(s,"Simulation: the data-generating process","Simulation setup")
bullets(s,[("n","Four visits, one binary severity covariate X. Randomized and external covariate mixes differ (shift, overlap preserved):",BODY,False)], 0.55,1.5,12.3,0.9, size=14.5)
img(s, E("eq_dgp_shift"), 0.7, 2.25, 9.6, 0.62, halign="left", bg=True)
bullets(s,[("n","Untreated control trajectory (random intercept + AR(1) errors), heteroscedastic scale, and the treatment effect:",BODY,False)], 0.55,3.0,12.3,0.7, size=14.5)
img(s, E("eq_dgp_y"),  0.7, 3.65, 11.4, 0.6, halign="left", bg=True)
img(s, E("eq_dgp_sig"),0.7, 4.35, 9.2, 0.55, halign="left", bg=True)
bullets(s,[("n","A fraction π_B of external controls are incompatible, with drift magnitude Δ (the signal-to-noise knob):",BODY,False)], 0.55,5.05,12.3,0.7, size=14.5)
img(s, E("eq_dgp_contam"),0.7, 5.7, 11.0, 0.6, halign="left", bg=True)
bullets(s,[
 ("n","Reference cell (everything is one axis away from it):  n₁=120 treated, n₀=60 RCT controls, 200 ECs, π_B=0.3, pattern A (global shift), Δ=8, estimand c₂, weighted CV+, correct model.",INK,True),
], 0.55,6.35,12.3,0.9, size=13.5)
page(s)

# --- Screen in action
s=newslide(); title(s,"The screen in action (one draw at the reference cell)","Simulation setup")
img(s, F("fig_outcome"), 6.7, 1.55, 6.2, 2.55, valign="top")
img(s, F("fig_pval"),    6.7, 4.25, 6.2, 2.55, valign="top")
bullets(s,[
 ("h","Top — raw outcome"),
 "Compatible ECs (blue) overlap the RCT controls (gray) but sit slightly higher, purely because the external mix carries more severe subjects (covariate shift, NOT incompatibility).",
 "Contaminated ECs (amber) sit far to the right. A raw-outcome threshold would confound shift with contamination.",
 ("h","Bottom — weighted residual-rank p-values"),
 "Contaminated ECs pile up near p = 0 and fall below γ — all screened out.",
 "Compatible ECs are roughly uniform on (0,1] and are kept; only about 6% fall below γ.",
 ("n","The screen truncates the calibrated p-value scale at γ, not the raw outcome.",INK,True),
], 0.55,1.5,6.0,5.5, size=13.5)
page(s)

# --- SNR
s=newslide(); title(s,"When can the screen detect contamination?","Simulation setup")
bullets(s,[
 "Detection is limited by the drift's signal-to-noise ratio Δ / σ:",
 ("s","Low Δ/σ — contamination is buried in the noise and leaks through ANY screen."),
 ("s","High Δ/σ — cleanly separable and reliably removed."),
 ("n","So borrowing is safest precisely when incompatibility is either absent, or large enough to catch.",INK,True),
 ("n","This also bounds the gain: the screen can only remove what it can see.",BODY,False),
], 0.55,1.55,12.2,2.3, size=15)
img(s, F("fig_snr"), 1.4, 3.9, 10.5, 3.1)
page(s)

# --- Results: reference cell
s=newslide(); title(s,"Result: beats RCT-only on all four criteria at once","Simulation results")
img(s, E("tab_ref"), 0.7, 1.6, 12.0, 2.55, bg=True)
bullets(s,[
 ("n","sym-ada (our default screen) is the only ATTAINABLE method that beats RCT-only on all four criteria simultaneously: unbiased (0.00), lower RMSE (0.18 vs 0.20), type-I no worse than RCT-only (0.06), higher power (0.73 vs 0.64).",INK,True),
 ("n","It reaches most of the way to the UNATTAINABLE oracle (0.78 power, 0.17 RMSE) — which knows the true labels.",BODY,False),
 ("n","Full borrowing is catastrophic: bias −0.42, type-I 0.42, coverage collapses.",RED,False),
 ("s","300 replicates, B = 500 full-pipeline bootstraps; read type-I relative to the RCT-only column (a shared bootstrap-CI floor)."),
], 0.55,4.4,12.3,2.6, size=14.5)
page(s)

# --- Results: robustness
s=newslide(); title(s,"Result: robust — and CV+ is the efficiency lever","Simulation results")
img(s, E("tab_robust"), 0.7, 1.6, 12.0, 2.35, bg=True)
bullets(s,[
 ("n","Robust to model misspecification: even dropping the covariate from the working mean model barely moves bias or power — the shared model error is absorbed by the RCT-control calibration.",BODY,False),
 ("n","Weighting is a low-cost safeguard: a near no-op under the correct model; it earns its keep under misspecification and at low signal-to-noise.",BODY,False),
 ("n","CV+ vs one-time split: +0.06 power (0.71 vs 0.65) at controlled type-I — simply by not wasting the scarce randomized controls.",INK,True),
], 0.55,4.2,12.3,2.8, size=14.5)
page(s)

# ============================================================ PART IV
section("Part IV", "Discussion", "Scope, limits, and open problems")

# --- Scope/limits
s=newslide(); title(s,"Scope, limits, and positioning","Discussion")
bullets(s,[
 ("n","Insurance, not a free lunch: residual-rank screening turns naive full-borrow's catastrophic bias into VALID, ROBUST inference, with a modest but real efficiency gain when the randomized arm is underpowered.",INK,True),
 ("h","The upside is bounded"),
 "An oracle-ceiling limits the gain to ≈ 15–25% RMSE: the treated arm is not borrowable, so the borrowable share of variance is  n(R,1) / (n(R,0)+n(R,1)).",
 "Borrowing helps most when randomized controls are scarce — but the screen also needs enough controls to calibrate; CV+ widens this window.",
 ("h","Two things you must choose well"),
 "The score must match the anticipated drift shape (a mismatched contrast is blind to the drift).",
 "Detection is signal-to-noise-limited; residual undetected low-SNR drift is the main risk.",
], 0.55,1.55,12.3,5.4, size=14.5)
page(s)

# --- Design questions: score choice + assumption hierarchy
s=newslide(); title(s,"Two design questions the theory raises","Discussion")
bullets(s,[
 ("h","1) Which nonconformity score?"),
 "Align the score to the anticipated drift: the estimand-contrast score for global / level shifts; the Mahalanobis (trajectory-shape) score for ramps, curvature, or transient dips; the max score for a single-visit spike.",
 "A score blind to the drift has no power — e.g. a visit-average contrast cannot see a shape change that averages out. Pick by the expected departure, and pre-specify it.",
 "Detection is not identification: a richer (VME-level) score catches more kinds of drift but can over-exclude compatible ECs; the estimand itself needs only ME-c.",
], 0.55,1.55,12.3,2.9, size=14.5)
bullets(s,[
 ("h","2) The assumption hierarchy — what is actually sufficient"),
 "Identifying the effect needs only ME-c (the one-contrast conditional mean); VME and DE are strictly stronger and are NOT required.",
 "Ordinary conformal validity nominally needs the full law DE; our screen relaxes this — it checks only the SCORE distribution (between ME-c and VME), a SUFFICIENT screen without assuming DE.",
 ("n","The crux is the gap between what identification NEEDS (ME-c), what the screen CHECKS (the score law), and what the chosen score is SENSITIVE to — pick the score so “score-compatible” implies “ME-c-compatible.”",INK,True),
], 0.55,4.5,12.3,2.55, size=13.5)
page(s)

# --- Open problems
s=newslide(); title(s,"Open problems","Discussion & problems")
bullets(s,[
 "Soft / continuous borrowing: down-weight rather than hard keep/drop, to remove the selection bias scalably (the current screen trades a little bias for robustness).",
 "Nested bootstrap for the adaptive-γ confidence interval: the current re-use of the point-selected γ inside the bootstrap is mildly anti-conservative.",
 "Type-I is reported RELATIVE to RCT-only — both share a percentile-bootstrap floor set by the number of resamples B; larger B pulls the absolute value to nominal.",
 "Multiplicity / FDR across many external controls (family-wise exclusion of at least one compatible EC is common).",
 "Validation on a covariate-driven (plasmode) design, beyond the analytic binary toy, to show the conclusions are not an artifact of the toy.",
], 0.55,1.55,12.3,5.4, size=15)
page(s)

# --- Summary
s=newslide(); title(s,"Summary","Takeaways")
bullets(s,[
 ("h","Problem"),
 "Borrow external controls to sharpen a hybrid trial — without importing bias from incompatible ones.",
 ("h","Idea"),
 "Rank each EC's longitudinal residual score among the randomized controls — weighted to absorb covariate shift — and borrow only the conforming ones through a transported AIPW.",
 ("h","Three design choices that make it work"),
 ("s","Symmetric trimming → the borrowed control mean stays unbiased."),
 ("s","CV+ rank calibration → does not waste the scarce randomized controls (the efficiency lever)."),
 ("s","Estimand-matched score → detects the drift that actually matters."),
 ("h","Why rank, not a raw residual"),
 ("n","Rank calibration is distribution-free and robust to model misspecification — the core reason. Result: safe, valid, robust borrowing that beats RCT-only and approaches the oracle.",INK,True),
], 0.55,1.5,12.3,5.6, size=14.5)
page(s)

out = os.path.join(HERE, "CC_presentation.pptx")
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
