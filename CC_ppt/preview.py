# Approximate visual preview of the .pptx: places the real embedded images at their
# true coordinates and wraps text the same way, so layout/overlap can be eyeballed.
import io, textwrap, sys
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

IN = 914400.0
prs = Presentation("CC_presentation.pptx")
SW, SH = prs.slide_width/IN, prs.slide_height/IN

def hexof(shape):
    try:
        c = shape.fill.fore_color.rgb
        return "#%02x%02x%02x" % (c[0], c[1], c[2])
    except Exception:
        return None

def runcol(p):
    for r in p.runs:
        try: c=r.font.color.rgb; return "#%02x%02x%02x"%(c[0],c[1],c[2])
        except Exception: pass
    return "#333333"

want = [int(x) for x in sys.argv[1:]] if len(sys.argv)>1 else list(range(1,len(prs.slides._sldIdLst)+1))
for i, sl in enumerate(prs.slides, 1):
    if i not in want: continue
    fig = plt.figure(figsize=(SW, SH), dpi=110)
    ax = fig.add_axes([0,0,1,1]); ax.set_xlim(0,SW); ax.set_ylim(0,SH); ax.axis("off")
    ax.add_patch(Rectangle((0,0),SW,SH,facecolor="white",edgecolor="#cccccc"))
    # pass 1: filled rectangles (backgrounds)
    for sh in sl.shapes:
        if sh.shape_type==MSO_SHAPE_TYPE.AUTO_SHAPE:
            hx=hexof(sh)
            if hx:
                l,t,w,h=sh.left/IN,sh.top/IN,sh.width/IN,sh.height/IN
                ax.add_patch(Rectangle((l,SH-t-h),w,h,facecolor=hx,edgecolor="none"))
    # pass 2: pictures
    for sh in sl.shapes:
        if sh.shape_type==MSO_SHAPE_TYPE.PICTURE:
            l,t,w,h=sh.left/IN,sh.top/IN,sh.width/IN,sh.height/IN
            try:
                im=Image.open(io.BytesIO(sh.image.blob)).convert("RGBA")
                ax.imshow(im, extent=[l,l+w,SH-t-h,SH-t], aspect="auto", zorder=3)
            except Exception:
                ax.add_patch(Rectangle((l,SH-t-h),w,h,facecolor="#dddddd",edgecolor="#999",zorder=3))
    # pass 3: text (manual wrap)
    for sh in sl.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip(): continue
        l,t,w=sh.left/IN,sh.top/IN,sh.width/IN
        y=SH-t-0.04
        for p in sh.text_frame.paragraphs:
            txt="".join(r.text for r in p.runs)
            if not txt.strip():
                y-=0.12; continue
            sizes=[r.font.size.pt for r in p.runs if r.font.size]; sz=max(sizes) if sizes else 14
            bold=any(r.font.bold for r in p.runs)
            col=runcol(p); indent=0.28 if (p.level or 0)>0 else 0.0
            cpl=max(6,int(w*118/sz))
            first=True
            for seg in txt.split("\n"):
                for line in (textwrap.wrap(seg,cpl) or [""]):
                    ax.text(l+indent, y, line, va="top", ha="left", fontsize=sz*0.98,
                            color=col, weight="bold" if bold else "normal", family="DejaVu Sans", zorder=4)
                    y-=1.28*sz/72.0
            y-=(p.space_after.pt if p.space_after is not None else 3)/72.0
    fig.savefig(f"/tmp/prev_s{i:02d}.png", dpi=110); plt.close(fig)
print("rendered", [f"/tmp/prev_s{i:02d}.png" for i in want[:1]], "...")
