import math
from pptx import Presentation
prs = Presentation("CC_presentation.pptx")
SWi, SHi = prs.slide_width/914400, prs.slide_height/914400
IN = 914400
for i, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        try:
            l,t,w,h = sh.left/IN, sh.top/IN, sh.width/IN, sh.height/IN
        except Exception:
            continue
        if l < -0.05 or t < -0.05 or l+w > SWi+0.08 or t+h > SHi+0.08:
            print(f"S{i:>2} OOB       ({l:5.2f},{t:5.2f}) {w:4.2f}x{h:4.2f}  {getattr(sh,'name','')[:18]}")
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf = sh.text_frame; th = 0.0
            for p in tf.paragraphs:
                txt = "".join(r.text for r in p.runs)
                sizes = [r.font.size.pt for r in p.runs if r.font.size]
                sz = max(sizes) if sizes else 14.0
                cpl = max(6.0, w*125.0/sz)                 # conservative chars/line
                nlines = max(1, math.ceil(len(txt)/cpl))
                th += nlines * 1.30 * sz/72.0
                th += (p.space_after.pt if p.space_after is not None else 0)/72.0
            if th > h*1.06 + 0.10:
                print(f"S{i:>2} OVERFLOW  est {th:4.2f} > box {h:4.2f}  '{tf.text[:44].strip()}'")
print("audit done; slides:", len(prs.slides._sldIdLst))
